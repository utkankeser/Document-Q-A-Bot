"""
rag_pipeline.py — RAG (Retrieval-Augmented Generation) Pipeline

Bu dosya Document Q&A Bot'un beynini oluşturur:
1. Dokümanlardan metin çıkarır (PDF, DOCX, TXT, PPT)
2. Metni parçalara (chunk) ayırır
3. Embedding oluşturur (sentence-transformers)
4. ChromaDB'ye kaydeder
5. Soru sorulduğunda en alakalı parçaları bulur
6. Gemini API'ye context + soru gönderir
"""

import uuid
import fitz  # PyMuPDF — PDF okumak için
from docx import Document as DocxDocument  # python-docx — DOCX okumak için
from pptx import Presentation  # python-pptx — PPT/PPTX okumak için
from sentence_transformers import SentenceTransformer
import chromadb
import google.generativeai as genai

from config import (
    GEMINI_API_KEY,
    CHROMA_DIR,
    CHUNK_SIZE,
    CHUNK_OVERLAP,
    TOP_K,
    EMBEDDING_MODEL,
    GEMINI_MODEL,
)

# ─── Model ve DB'yi başlat ─────────────────────────────────────
# Embedding modeli — metni vektöre çevirir
# İlk çalışmada modeli indirir (~80MB), sonraki çalışmalarda cache'den kullanır
embedding_model = SentenceTransformer(EMBEDDING_MODEL)

# ChromaDB — vektörleri sakladığımız veritabanı
chroma_client = chromadb.PersistentClient(path=str(CHROMA_DIR))
collection = chroma_client.get_or_create_collection(
    name="documents",
    metadata={"hnsw:space": "cosine"},  # Cosine similarity kullan
)

# Gemini API'yi yapılandır
genai.configure(api_key=GEMINI_API_KEY)


# ─── 1. Metin Çıkarma ──────────────────────────────────────────

def extract_text_from_pdf(file_path: str) -> str:
    """PDF dosyasından metin çıkarır."""
    doc = fitz.open(file_path)
    text = ""
    for page in doc:
        text += page.get_text()
    doc.close()
    return text


def extract_text_from_docx(file_path: str) -> str:
    """DOCX dosyasından metin çıkarır."""
    doc = DocxDocument(file_path)
    text = "\n".join([para.text for para in doc.paragraphs])
    return text


def extract_text_from_txt(file_path: str) -> str:
    """TXT dosyasından metin okur."""
    with open(file_path, "r", encoding="utf-8") as f:
        return f.read()


def extract_text_from_pptx(file_path: str) -> str:
    """PPT/PPTX dosyasından metin çıkarır."""
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text += paragraph.text + "\n"
    return text


def extract_text(file_path: str) -> str:
    """
    Dosya uzantısına göre uygun metin çıkarma fonksiyonunu çağırır.

    Desteklenen formatlar: PDF, DOCX, TXT, PPT/PPTX
    """
    ext = file_path.lower().rsplit(".", 1)[-1]
    extractors = {
        "pdf": extract_text_from_pdf,
        "docx": extract_text_from_docx,
        "txt": extract_text_from_txt,
        "pptx": extract_text_from_pptx,
        "ppt": extract_text_from_pptx,
    }
    extractor = extractors.get(ext)
    if not extractor:
        raise ValueError(f"Desteklenmeyen dosya formatı: .{ext}")
    return extractor(file_path)


# ─── 2. Metin Parçalama (Chunking) ─────────────────────────────

def chunk_text(text: str) -> list[str]:
    """
    Metni daha küçük parçalara (chunk) ayırır.

    Neden parçalıyoruz?
    - Embedding modelleri uzun metinleri iyi işleyemez
    - Daha küçük parçalar daha hassas arama sağlar
    - LLM'in context window'una sığması gerekir

    Overlap: Ardışık parçalar arasında örtüşme bırakıyoruz
    ki bir cümle parçalar arasında bölünürse bağlam kaybolmasın.
    """
    if not text.strip():
        return []

    chunks = []
    start = 0
    while start < len(text):
        end = start + CHUNK_SIZE
        chunk = text[start:end]
        if chunk.strip():  # Boş chunk ekleme
            chunks.append(chunk.strip())
        start += CHUNK_SIZE - CHUNK_OVERLAP  # Overlap kadar geri gel
    return chunks


# ─── 3. Doküman İşleme (Embedding + Kaydetme) ──────────────────

def process_document(file_path: str, doc_id: str, filename: str) -> dict:
    """
    Bir dokümanı uçtan uca işler:
    1. Metni çıkar
    2. Parçalara ayır
    3. Her parçanın embedding'ini oluştur
    4. ChromaDB'ye kaydet

    Returns:
        dict: İşlem sonucu bilgileri (chunk sayısı vb.)
    """
    # Metin çıkar
    text = extract_text(file_path)
    if not text.strip():
        raise ValueError("Doküman boş veya metin çıkarılamadı.")

    # Parçalara ayır
    chunks = chunk_text(text)
    if not chunks:
        raise ValueError("Doküman parçalanamadı.")

    # Embedding oluştur — her chunk için bir vektör
    embeddings = embedding_model.encode(chunks).tolist()

    # ChromaDB'ye kaydet
    ids = [f"{doc_id}_chunk_{i}" for i in range(len(chunks))]
    metadatas = [{"doc_id": doc_id, "filename": filename, "chunk_index": i} for i in range(len(chunks))]

    collection.add(
        ids=ids,
        embeddings=embeddings,
        documents=chunks,
        metadatas=metadatas,
    )

    return {
        "doc_id": doc_id,
        "filename": filename,
        "total_chunks": len(chunks),
        "text_length": len(text),
    }


# ─── 4. Soru Sorma (Retrieval + Generation) ────────────────────

def retrieve_context(query: str, doc_id: str | None = None) -> list[str]:
    """
    Soruya en yakın metin parçalarını ChromaDB'den getirir.

    1. Soruyu embedding'e çevir
    2. ChromaDB'de en yakın vektörleri bul (cosine similarity)
    3. İlgili metin parçalarını döndür
    """
    query_embedding = embedding_model.encode([query]).tolist()

    # Belirli bir dokümana filtre uygula (isteğe bağlı)
    where_filter = {"doc_id": doc_id} if doc_id else None

    results = collection.query(
        query_embeddings=query_embedding,
        n_results=TOP_K,
        where=where_filter,
    )

    # Sonuçları döndür
    if results and results["documents"]:
        return results["documents"][0]
    return []


def ask_question(query: str, doc_id: str | None = None) -> dict:
    """
    Kullanıcının sorusunu cevaplar:
    1. İlgili context'i ChromaDB'den getir
    2. Gemini API'ye context + soru gönder
    3. Cevabı döndür
    """
    # İlgili parçaları getir
    context_chunks = retrieve_context(query, doc_id)

    if not context_chunks:
        return {
            "answer": "Bu soru için dokümanlarınızda ilgili bir bilgi bulamadım. "
                      "Lütfen farklı bir soru deneyin veya yeni bir doküman yükleyin.",
            "context_used": [],
        }

    # Context'i birleştir
    context = "\n\n---\n\n".join(context_chunks)

    # Gemini'ye gönderilecek prompt
    prompt = f"""Sen bir doküman asistanısın. Aşağıdaki doküman parçalarına dayanarak kullanıcının sorusunu cevapla.

ÖNEMLİ KURALLAR:
- SADECE verilen doküman parçalarındaki bilgilere dayanarak cevap ver.
- Bilgi dokümanlardan gelmiyorsa, "Bu bilgi dokümanlarınızda bulunamadı" de.
- Cevabın açık, net ve anlaşılır olsun.
- Gerekirse maddeler halinde listele.

--- DOKÜMAN PARÇALARI ---
{context}
--- DOKÜMAN PARÇALARI SONU ---

KULLANICININ SORUSU: {query}

CEVAP:"""

    # Gemini API'ye gönder
    model = genai.GenerativeModel(GEMINI_MODEL)
    response = model.generate_content(prompt)

    return {
        "answer": response.text,
        "context_used": context_chunks,
    }


# ─── 5. Doküman Silme ──────────────────────────────────────────

def delete_document(doc_id: str) -> bool:
    """Bir dokümanın tüm chunk'larını ChromaDB'den siler."""
    try:
        # Bu doc_id'ye ait tüm kayıtları sil
        results = collection.get(where={"doc_id": doc_id})
        if results and results["ids"]:
            collection.delete(ids=results["ids"])
        return True
    except Exception:
        return False


def get_all_documents() -> list[dict]:
    """ChromaDB'deki tüm dokümanları (benzersiz doc_id'leri) döndürür."""
    try:
        all_data = collection.get()
        if not all_data or not all_data["metadatas"]:
            return []

        # Benzersiz dokümanları topla
        docs = {}
        for meta in all_data["metadatas"]:
            did = meta["doc_id"]
            if did not in docs:
                docs[did] = {
                    "doc_id": did,
                    "filename": meta["filename"],
                    "chunk_count": 0,
                }
            docs[did]["chunk_count"] += 1

        return list(docs.values())
    except Exception:
        return []
